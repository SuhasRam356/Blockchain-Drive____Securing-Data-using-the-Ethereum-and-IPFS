import collections 
import collections.abc
import sys

# python-pptx needs collections.abc.Container which was removed in Python 3.10
# We monkey-patch it here so python-pptx doesn't crash on Python 3.13
if not hasattr(collections, 'Container'):
    collections.Container = collections.abc.Container

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

BG_IMAGE = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\presentation_bg_light_1785668137706.png"

def apply_background(slide, prs):
    if os.path.exists(BG_IMAGE):
        pic = slide.shapes.add_picture(BG_IMAGE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        # Send picture to back by moving its XML element
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

def add_slide(prs, title_text, bullets, notes_text=None, image_path=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    apply_background(slide, prs)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # Dark IEEE blue text for light BG
    
    body_shape = slide.shapes.placeholders[1]
    
    if image_path and os.path.exists(image_path):
        body_shape.width = Inches(5.0)
        img_left = Inches(5.5)
        img_top = Inches(2.0)
        img_width = Inches(4.0)
        try:
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
        except Exception as e:
            print(f"Error adding image: {e}")
            
    tf = body_shape.text_frame
    tf.clear()
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(30, 30, 30) # Dark gray for bullets
        
    if notes_text and slide.has_notes_slide:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

def add_title_slide(prs, title_text, subtitle_text, bullets, notes_text=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    apply_background(slide, prs)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # Dark blue for title
    
    subtitle.text = subtitle_text + "\n\n" + "\n".join(bullets)
    for p in subtitle.text_frame.paragraphs:
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(30, 30, 30) # Dark gray
    
    if notes_text and slide.has_notes_slide:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

def main():
    prs = Presentation()
    
    # Image Paths (from AI generation)
    img_title = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\img_title_1785668483126.png"
    img_intro = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\img_intro_1785668494039.png"
    img_limitations = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\img_limitations_1785668503677.png"
    img_solution = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\img_solution_1785668515656.png"
    img_tech = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\img_tech_1785668536075.png"
    img_ipfs = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\img_ipfs_1785668546577.png"
    img_smart_contract = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\img_smart_contract_1785668559216.png"
    img_upload = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\img_upload_1785668570539.png"
    img_download = r"C:\Users\suhas\.gemini\antigravity-ide\brain\0e043bbe-b2fe-41af-8b3a-bdef9a6feb5f\img_download_1785668590608.png"

    add_title_slide(prs, 
        "Blockchain Drive: Securing Data using Ethereum and IPFS", 
        "A Decentralized Approach to Cloud Storage System Architecture", 
        ["Project Presentation", "Your Name", "Project Guide's Name", "Date"], 
        "Good morning/afternoon everyone. Today, I am excited to present my project, 'Blockchain Drive', a secure, decentralized alternative to traditional cloud storage like Google Drive or Dropbox.")
    
    add_slide(prs, "I. Introduction: The Era of Cloud Storage", 
        ["Massive Data Growth: Generating millions of gigabytes daily.",
         "Reliance on Cloud: Entrusting 3rd-party servers with sensitive data.",
         "The Fundamental Flaw: Users relinquish absolute ownership of their data."],
        "We live in a digital world where cloud storage is essential.", img_intro)
    
    add_slide(prs, "II. Background: Centralized Storage Limitations",
        ["Single Point of Failure (SPoF): Server outages completely deny access.",
         "Data Breaches: Centralized repositories are high-value targets for malicious actors.",
         "Censorship & Provider Lock-in: Unilateral account suspension or arbitrary policy changes."],
        "Because current systems are centralized, they have a single point of failure.", img_limitations)
        
    add_slide(prs, "III. Proposed Solution: Decentralized Architecture",
        ["Blockchain Drive: A Web3-powered file storage paradigm.",
         "Distributed Network: Mitigation of SPoF via global node distribution.",
         "Cryptographic Ownership: Users hold absolute private keys.",
         "Tamper-Proof Logic: Execution governed by Ethereum Smart Contracts."],
        "To solve these issues, I built Blockchain Drive.", img_solution)
        
    add_slide(prs, "IV. System Technology Stack",
        ["Frontend Interface: React.js, TailwindCSS, Ethers.js",
         "Smart Contract Layer: Solidity, Hardhat, OpenZeppelin",
         "Storage Layer: IPFS (InterPlanetary File System), Pinata",
         "Consensus Network: Ethereum (Sepolia Testnet)"],
        "Building this required a modern Web3 tech stack.", img_tech)
        
    add_slide(prs, "V. Core Component 1: InterPlanetary File System (IPFS)",
        ["Peer-to-Peer Protocol: Eradicates centralized server dependency.",
         "Content-Addressed Storage: Retrieval via cryptographic hash (CID) rather than location.",
         "Immutability: Any file modification inherently changes its cryptographic CID."],
        "The backbone of our storage is IPFS.", img_ipfs)
        
    add_slide(prs, "VI. Core Component 2: Ethereum Smart Contracts",
        ["Immutable State Ledger: Cryptographically records file-to-owner mapping.",
         "Deterministic Execution: Code operates exactly as deployed without intermediaries.",
         "Decentralized Access Control: Manages read/write permissions at the protocol level."],
        "While IPFS stores the file, the Ethereum blockchain stores the permissions.", img_smart_contract)
        
    add_slide(prs, "VII. System Workflow: Data Ingestion (Upload)",
        ["1. User initiates file selection via frontend client.",
         "2. File undergoes Client-Side AES-256 Encryption.",
         "3. Ciphertext is transmitted to the IPFS decentralized network.",
         "4. IPFS generates and returns a unique Content Identifier (CID).",
         "5. CID and access parameters are committed to the Ethereum Smart Contract."],
        "Let's walk through how uploading works.", img_upload)
        
    add_slide(prs, "VIII. System Workflow: Data Retrieval (Download)",
        ["1. Client queries Smart Contract for authorized CIDs.",
         "2. Smart Contract validates cryptographic signature of the requester.",
         "3. Client retrieves ciphertext from IPFS using the authorized CID.",
         "4. Client-side decryption restores plaintext data."],
        "Downloading is just the reverse.", img_download)
        
    add_slide(prs, "IX. Security Implementation: Cryptography",
        ["End-to-End Encryption (E2EE): Absolute data confidentiality.",
         "Symmetric Encryption (AES-256): Standardized, robust ciphertext generation.",
         "Zero-Knowledge Architecture: Decentralized nodes store routing data, entirely blind to plaintext payloads."],
        "Security is our highest priority.", img_tech)
        
    add_slide(prs, "X. IEEE Feature 1: Convergent Encryption (Deduplication)",
        ["The Problem: Independent encryption of identical files causes massive storage redundancy.",
         "The Solution: Convergent Encryption utilizes a cryptographic hash of the plaintext as the encryption key.",
         "Result: Cross-user Data Deduplication achieved alongside zero-knowledge privacy."],
        "To make this project stand out, I implemented Convergent Encryption.", img_ipfs)
        
    add_slide(prs, "XI. Distributed Access Control",
        ["Cryptographic Sharing: Key distribution via asymmetric encryption (simulated).",
         "On-Chain Authorization: Smart contract acts as the access control list (ACL).",
         "Granular Control: Absolute owner sovereignty over decryption capabilities."],
        "Our system also allows secure file sharing.", img_solution)
        
    add_slide(prs, "XII. IEEE Feature 2: Simulated Proxy Re-Encryption (PRE)",
        ["The Problem: Access revocation is ineffective if the revoked user retained the decryption key.",
         "The Solution: Key Rotation mechanism simulating PRE.",
         "Result: Revocation automatically triggers global key rotation, cryptographically excluding the revoked user."],
        "Another major IEEE feature I added is Key Rotation.", img_limitations)
        
    add_slide(prs, "XIII. Enhancing Usability: Decentralized Identity (ENS)",
        ["The Problem: Hexadecimal Ethereum addresses are prone to human error.",
         "The Solution: Ethereum Name Service (ENS) integration.",
         "Impact: Replaces complex addresses with human-readable namespaces (e.g., alice.eth)."],
        "To improve the user experience, I integrated the Ethereum Name Service.", img_title)
        
    add_slide(prs, "XIV. Protocol Governance: DAO Integration",
        ["Decentralized Autonomous Organization (DAO): Governance by stakeholders.",
         "DriveToken (DRV): Sybil-resistant voting weight mechanism.",
         "Decentralized Upgrades: Protocol changes execute strictly via on-chain community consensus."],
        "Finally, this platform is governed by a DAO.", img_intro)
        
    add_slide(prs, "XV. Conclusion and Future Directions",
        ["Achieved Objectives: Zero downtime, Absolute Data Sovereignty, Cryptographic Privacy.",
         "Future Work 1: Layer-2 (L2) Rollup integration for gas cost mitigation.",
         "Future Work 2: Advanced Zero-Knowledge Proofs (zk-SNARKs) for anonymous access verification."],
        "In conclusion, Blockchain Drive offers total data sovereignty and zero downtime.", img_tech)
        
    add_slide(prs, "XVI. Thank You",
        ["Questions & Discussion",
         "Live System Demonstration",
         "Project Repository Available"],
        "Thank you all for your time and attention.")
        
    prs.save('Blockchain_Drive_IEEE_Presentation_v4.pptx')
    print("Successfully created Blockchain_Drive_IEEE_Presentation_v4.pptx")

if __name__ == "__main__":
    main()
